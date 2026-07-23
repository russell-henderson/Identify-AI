from __future__ import annotations

import asyncio
import inspect
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Awaitable, Callable

import pymupdf as fitz
from ollama import AsyncClient
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from utils import load_config


StatusCallback = Callable[[str], Awaitable[None] | None]
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
WINDOWS_RESERVED_NAMES = {
    "CON",
    "PRN",
    "AUX",
    "NUL",
    *(f"COM{number}" for number in range(1, 10)),
    *(f"LPT{number}" for number in range(1, 10)),
}

RESULT_SCHEMA: dict[str, Any] = {
    "type": "object",
    "properties": {
        "category": {"type": "string"},
        "suggested_name": {"type": "string"},
        "reasoning": {"type": "string"},
    },
    "required": ["category", "suggested_name", "reasoning"],
    "additionalProperties": False,
}


async def check_ollama(config: dict[str, Any]) -> str | None:
    """Return a user-facing preflight error, or None when Ollama is ready."""
    try:
        ollama_config = config["ollama"]
        client = AsyncClient(host=ollama_config.get("host", "http://127.0.0.1:11434"), timeout=10.0)
        response = await client.list()
        models = getattr(response, "models", None)
        if models is None and isinstance(response, dict):
            models = response.get("models", [])
        installed = {_model_name(model) for model in (models or [])}
        required = {str(ollama_config["vision_model"]), str(ollama_config["text_model"])}
        missing = sorted(required.difference(installed))
        if missing:
            return "Missing Ollama model(s): " + ", ".join(missing) + ". Run ollama pull for each model."
        return None
    except Exception as exc:
        return _friendly_error(exc)


def _model_name(model: Any) -> str:
    if isinstance(model, dict):
        return str(model.get("model") or model.get("name") or "")
    return str(getattr(model, "model", None) or getattr(model, "name", None) or "")


async def analyze_file(file_path: Path, callback: StatusCallback) -> dict[str, str]:
    """
    Analyze one local file with Ollama and return a normalized naming proposal.

    Errors are converted into an error result so one bad file cannot terminate a
    larger batch.
    """
    path = Path(file_path)

    try:
        config = await load_config()
        await _emit(callback, f"Inspecting {path.name}...")

        if not path.exists():
            raise FileNotFoundError(f"File does not exist: {path}")
        if not path.is_file():
            raise ValueError(f"Path is not a file: {path}")

        extension = path.suffix.lower()
        ollama_config = config["ollama"]
        client = AsyncClient(
            host=ollama_config.get("host", "http://127.0.0.1:11434"),
            timeout=600.0,
        )

        prompt = await _build_prompt(path, config)

        if extension in IMAGE_EXTENSIONS:
            model = str(ollama_config["vision_model"])
            image_bytes = await _load_visual_content(path, callback, config)
            await _emit(callback, f"Analyzing visual content with {model}...")
            messages = [
                {
                    "role": "user",
                    "content": prompt,
                    "images": [image_bytes],
                }
            ]
        elif extension == ".pdf":
            await _emit(callback, "Extracting PDF text...")
            text_content = await _read_pdf_text(path, config)
            if text_content:
                model = str(ollama_config["text_model"])
                await _emit(callback, f"Analyzing PDF text with {model}...")
                messages = [{"role": "user", "content": f"{prompt}\n\nPDF TEXT:\n{text_content}"}]
            else:
                model = str(ollama_config["vision_model"])
                image_bytes = await _load_visual_content(path, callback, config)
                await _emit(callback, f"Analyzing scanned PDF with {model}...")
                messages = [{"role": "user", "content": prompt, "images": [image_bytes]}]
        else:
            model = str(ollama_config["text_model"])
            await _emit(callback, "Reading text content...")
            text_content = await _read_text_content(path, config)
            await _emit(callback, f"Analyzing text with {model}...")
            messages = [
                {
                    "role": "user",
                    "content": f"{prompt}\n\nFILE CONTENT:\n{text_content}",
                }
            ]

        response = await client.chat(
            model=model,
            messages=messages,
            format=RESULT_SCHEMA,
            stream=False,
            options={
                "temperature": float(ollama_config.get("temperature", 0.1)),
                "num_ctx": int(ollama_config.get("context_length", 8192)),
            },
            keep_alive=ollama_config.get("keep_alive", "10m"),
        )

        raw_content = response.message.content or ""
        await _emit(callback, "Validating model response...")
        parsed = _parse_model_json(raw_content)
        normalized = await _normalize_result(parsed, path, config)
        await _emit(callback, "Analysis complete")
        return normalized

    except Exception as exc:
        error_message = _friendly_error(exc)
        await _emit(callback, f"Error: {error_message}")
        return {
            "category": "error",
            "suggested_name": path.name,
            "reasoning": error_message,
        }


async def _emit(callback: StatusCallback, message: str) -> None:
    try:
        result = callback(message)
        if inspect.isawaitable(result):
            await result
    except Exception:
        # UI callback failures must not abort the file analysis itself.
        return


async def _load_visual_content(
    file_path: Path,
    callback: StatusCallback,
    config: dict[str, Any],
) -> bytes:
    if file_path.suffix.lower() == ".pdf":
        await _emit(callback, "Rendering first PDF page...")
        dpi = int(config["processing"].get("pdf_render_dpi", 150))
        return await asyncio.to_thread(_render_first_pdf_page, file_path, dpi)

    await _emit(callback, "Reading image bytes...")
    return await asyncio.to_thread(file_path.read_bytes)


def _render_first_pdf_page(file_path: Path, dpi: int) -> bytes:
    with fitz.open(file_path) as document:
        if document.page_count < 1:
            raise ValueError("PDF contains no pages")
        if document.needs_pass:
            raise PermissionError("PDF is password protected")

        page = document.load_page(0)
        pixmap = page.get_pixmap(
            dpi=max(72, min(dpi, 300)),
            colorspace=fitz.csRGB,
            alpha=False,
        )
        return pixmap.tobytes("png")


async def _read_text_content(file_path: Path, config: dict[str, Any]) -> str:
    max_characters = int(config["processing"].get("max_text_characters", 24000))
    return await asyncio.to_thread(_read_text_content_sync, file_path, max_characters)


def _read_text_content_sync(file_path: Path, max_characters: int) -> str:
    extension = file_path.suffix.lower()
    if extension == ".docx":
        return _limit_text(_extract_docx(file_path), max_characters)
    if extension == ".xlsx":
        return _limit_text(_extract_xlsx(file_path), max_characters)
    if extension == ".pptx":
        return _limit_text(_extract_pptx(file_path), max_characters)

    byte_limit = max(4096, max_characters * 4)
    with file_path.open("rb") as source_file:
        raw = source_file.read(byte_limit + 1)

    truncated = len(raw) > byte_limit
    raw = raw[:byte_limit]

    if raw and raw.count(b"\x00") / len(raw) > 0.02:
        raise ValueError(
            "The file appears to be binary. Add a dedicated extractor before "
            "processing this file type as text."
        )

    text = _decode_text(raw)
    text = text[:max_characters]

    if not text.strip():
        raise ValueError("The file contains no readable text")

    if truncated or len(text) >= max_characters:
        text += "\n\n[CONTENT TRUNCATED FOR LOCAL MODEL CONTEXT LIMIT]"

    return text


async def _read_pdf_text(file_path: Path, config: dict[str, Any]) -> str:
    max_characters = int(config["processing"].get("max_text_characters", 24000))
    return await asyncio.to_thread(_read_pdf_text_sync, file_path, max_characters)


def _read_pdf_text_sync(file_path: Path, max_characters: int) -> str:
    with fitz.open(file_path) as document:
        if document.page_count < 1:
            raise ValueError("PDF contains no pages")
        if document.needs_pass:
            raise PermissionError("PDF is password protected")
        parts = []
        for index, page in enumerate(document, start=1):
            page_text = page.get_text("text").strip()
            if page_text:
                parts.append(f"[Page {index}]\n{page_text}")
    text = "\n\n".join(parts)
    # Scanned PDFs often have no usable text; the caller will use the vision model.
    return _limit_text(text, max_characters, allow_empty=True)


def _extract_docx(file_path: Path) -> str:
    document = Document(file_path)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    tables = [
        " | ".join(cell.text.strip() for cell in row.cells)
        for table in document.tables
        for row in table.rows
        if any(cell.text.strip() for cell in row.cells)
    ]
    return "\n".join(paragraphs + tables)


def _extract_xlsx(file_path: Path) -> str:
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    try:
        lines: list[str] = []
        for worksheet in workbook.worksheets:
            lines.append(f"[Sheet: {worksheet.title}]")
            for row in worksheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value not in (None, "")]
                if values:
                    lines.append(" | ".join(values))
        return "\n".join(lines)
    finally:
        workbook.close()


def _extract_pptx(file_path: Path) -> str:
    presentation = Presentation(file_path)
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[Slide {index}]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def _limit_text(text: str, max_characters: int, allow_empty: bool = False) -> str:
    normalized = text.strip()
    if not normalized:
        if allow_empty:
            return ""
        raise ValueError("The file contains no readable text")
    if len(normalized) > max_characters:
        return normalized[:max_characters] + "\n\n[CONTENT TRUNCATED FOR LOCAL MODEL CONTEXT LIMIT]"
    return normalized


def _decode_text(raw: bytes) -> str:
    encodings = ("utf-8-sig", "utf-16", "utf-16-le", "utf-16-be", "cp1252")
    for encoding in encodings:
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


async def _build_prompt(file_path: Path, config: dict[str, Any]) -> str:
    stat_result = file_path.stat()
    modified_date = datetime.fromtimestamp(stat_result.st_mtime).strftime(
        config["naming"].get("date_format", "%Y-%m-%d")
    )
    naming = config["naming"]
    active_schema = naming.get("active_schema", "default")
    schema_template = naming.get("schemas", {}).get(
        active_schema,
        "{date}_{category}_{title}{extension}",
    )

    return f"""
You are a local document-classification and file-naming engine.
Analyze the supplied file and identify its most useful subject and document type.

Return strictly one JSON object and no commentary, markdown, or code fences.
The JSON must contain exactly these string keys:
- category: a short stable category such as invoice, receipt, contract, photo,
  report, notes, source_code, dataset, correspondence, or reference.
- suggested_name: a concise descriptive title only. Do not include a directory,
  date prefix, category prefix, or file extension.
- reasoning: one or two sentences explaining the classification and title.

Naming context:
- Original filename: {file_path.name}
- Original extension: {file_path.suffix.lower() or "none"}
- File modified date: {modified_date}
- Application naming schema: {schema_template}

Avoid vague titles such as document, image, scan, file, untitled, or miscellaneous.
Do not invent names, account numbers, dates, organizations, or events that are not
visible in the supplied content.
""".strip()


def _strip_markdown_code_fences(content: str) -> str:
    cleaned = content.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    return cleaned.strip()


def _parse_model_json(content: str) -> dict[str, str]:
    cleaned = _strip_markdown_code_fences(content)

    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        start = cleaned.find("{")
        end = cleaned.rfind("}")
        if start == -1 or end == -1 or end <= start:
            raise ValueError("Ollama did not return a JSON object")
        parsed = json.loads(cleaned[start : end + 1])

    if not isinstance(parsed, dict):
        raise ValueError("Ollama returned JSON, but it was not an object")

    missing = {"category", "suggested_name", "reasoning"}.difference(parsed)
    if missing:
        raise ValueError(
            "Ollama response is missing required keys: " + ", ".join(sorted(missing))
        )

    return {
        "category": str(parsed["category"]).strip(),
        "suggested_name": str(parsed["suggested_name"]).strip(),
        "reasoning": str(parsed["reasoning"]).strip(),
    }


async def _normalize_result(
    result: dict[str, str],
    file_path: Path,
    config: dict[str, Any],
) -> dict[str, str]:
    stat_result = file_path.stat()
    naming = config["naming"]
    separator = str(naming.get("separator", "_")) or "_"
    lowercase = bool(naming.get("lowercase", False))
    max_stem_length = max(24, int(naming.get("max_stem_length", 110)))

    category = _sanitize_component(result["category"], separator, lowercase)
    title = _sanitize_component(result["suggested_name"], separator, lowercase)

    if not category:
        category = "uncategorized"
    if not title:
        title = _sanitize_component(file_path.stem, separator, lowercase) or "untitled"

    extension = file_path.suffix.lower()
    modified_date = datetime.fromtimestamp(stat_result.st_mtime).strftime(
        naming.get("date_format", "%Y-%m-%d")
    )
    schema_name = naming.get("active_schema", "default")
    schema = naming.get("schemas", {}).get(
        schema_name,
        "{date}_{category}_{title}{extension}",
    )

    try:
        candidate = schema.format(
            date=modified_date,
            category=category,
            title=title,
            extension=extension,
            original_stem=_sanitize_component(file_path.stem, separator, lowercase),
        )
    except (KeyError, ValueError) as exc:
        raise ValueError(f"Invalid naming schema '{schema_name}': {exc}") from exc

    suggested_filename = normalize_user_filename(candidate, file_path, config)
    return {
        "category": category,
        "suggested_name": suggested_filename,
        "reasoning": result["reasoning"] or "Classified from the file content.",
    }


def normalize_user_filename(value: str, file_path: Path, config: dict[str, Any]) -> str:
    """Return an editable filename normalized to the source file's extension."""
    naming = config["naming"]
    separator = str(naming.get("separator", "_")) or "_"
    lowercase = bool(naming.get("lowercase", False))
    max_stem_length = max(24, int(naming.get("max_stem_length", 110)))
    candidate_stem = _sanitize_filename_stem(Path(value).stem, separator, lowercase)
    candidate_stem = candidate_stem[:max_stem_length].rstrip(" ._-" + separator)
    if not candidate_stem:
        candidate_stem = "untitled"
    if candidate_stem.upper() in WINDOWS_RESERVED_NAMES:
        candidate_stem = f"_{candidate_stem}"
    return f"{candidate_stem}{file_path.suffix.lower()}"


def _sanitize_component(value: str, separator: str, lowercase: bool) -> str:
    value = str(value).strip()
    value = re.sub(r"[<>:\"/\\|?*\x00-\x1F]", " ", value)
    value = re.sub(r"[\s._-]+", separator, value)
    value = value.strip(f" {separator}.")
    return value.lower() if lowercase else value


def _sanitize_filename_stem(value: str, separator: str, lowercase: bool) -> str:
    value = re.sub(r"[<>:\"/\\|?*\x00-\x1F]", " ", str(value))
    value = re.sub(r"\s+", separator, value)
    value = re.sub(rf"{re.escape(separator)}{{2,}}", separator, value)
    value = value.strip(f" {separator}.")
    return value.lower() if lowercase else value


def _friendly_error(exc: Exception) -> str:
    message = str(exc).strip() or exc.__class__.__name__
    lowered = message.lower()

    if "connection" in lowered and "refused" in lowered:
        return "Cannot connect to Ollama at the configured host. Start Ollama and retry."
    if "not found" in lowered and "model" in lowered:
        return f"Required Ollama model is not installed: {message}"
    if isinstance(exc, json.JSONDecodeError):
        return f"The model returned malformed JSON: {message}"

    return message
