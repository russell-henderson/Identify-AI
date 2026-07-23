import json
import tempfile
import unittest
from pathlib import Path

from docx import Document
from openpyxl import Workbook
import pymupdf as fitz
from pptx import Presentation

from analyzer import _model_name, _read_pdf_text_sync, _read_text_content_sync, normalize_user_filename
from main import build_file_row, scan_folder, write_batch_artifacts


class MvpWorkflowTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)
        self.config = {
            "ollama": {"vision_model": "vision", "text_model": "text"},
            "naming": {"separator": "_", "max_stem_length": 110, "lowercase": False},
        }

    def tearDown(self):
        self.temp_dir.cleanup()

    def test_user_filename_is_safe_and_keeps_source_extension(self):
        result = normalize_user_filename("CON: quarterly report.xlsx", self.root / "source.pdf", self.config)
        self.assertEqual(result, "CON_quarterly_report.pdf")

    def test_model_name_supports_dict_and_client_objects(self):
        class Model:
            model = "local-model:latest"

        self.assertEqual(_model_name({"name": "dict-model:latest"}), "dict-model:latest")
        self.assertEqual(_model_name(Model()), "local-model:latest")

    def test_scan_folder_respects_recursion_and_extensions(self):
        (self.root / "nested").mkdir()
        (self.root / "one.pdf").write_text("one")
        (self.root / "nested" / "two.docx").write_text("not a real document")
        self.assertEqual([path.name for path in scan_folder(self.root, {".pdf", ".docx"}, False)], ["one.pdf"])
        self.assertEqual([path.name for path in scan_folder(self.root, {".pdf", ".docx"}, True)], ["two.docx", "one.pdf"])

    def test_office_text_extractors_read_content(self):
        docx_path = self.root / "brief.docx"
        document = Document()
        document.add_paragraph("Project Aurora brief")
        document.save(docx_path)

        xlsx_path = self.root / "budget.xlsx"
        workbook = Workbook()
        workbook.active.append(["Department", "Budget"])
        workbook.active.append(["Research", 100])
        workbook.save(xlsx_path)
        workbook.close()

        pptx_path = self.root / "review.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[1]).shapes.title.text = "Quarterly Review"
        presentation.save(pptx_path)

        self.assertIn("Project Aurora", _read_text_content_sync(docx_path, 1000))
        self.assertIn("Research | 100", _read_text_content_sync(xlsx_path, 1000))
        self.assertIn("Quarterly Review", _read_text_content_sync(pptx_path, 1000))

    def test_scanned_pdf_without_text_falls_back_to_vision_path(self):
        pdf_path = self.root / "scan.pdf"
        document = fitz.open()
        document.new_page()
        document.save(pdf_path)
        document.close()
        self.assertEqual(_read_pdf_text_sync(pdf_path, 1000), "")

    def test_artifacts_record_staged_copy(self):
        source = self.root / "source.pdf"
        source.write_bytes(b"placeholder")
        row = build_file_row(source, "")
        row.proposed_name.value = "2026_report.pdf"
        row.category.value = "report"
        row.status.value = "Staged: 2026_report.pdf"
        row.approved.value = True
        row.staged_path = self.root / "Staging" / "2026_report.pdf"
        report, manifest = write_batch_artifacts([source], [row], None, self.root / "Archive", self.config, "2026-01-02_03-04-05")
        self.assertTrue(report.exists())
        payload = json.loads(manifest.read_text(encoding="utf-8"))
        self.assertEqual(payload["files"][0]["staged_path"], str(row.staged_path))


if __name__ == "__main__":
    unittest.main()
